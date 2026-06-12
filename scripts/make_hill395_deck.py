#!/usr/bin/env python
"""Build the CITA-Net / Hill 395 presentation deck (.pptx).

Self-contained: pulls the real metrics + dataset stats from
hill395_experiment_results/ and embeds the generated figures. Formulas mirror the
actual implementation (cta.py, encoder.py, graph_encoder.py, decoder.py,
losses.py, kinematics.py).
"""
from __future__ import annotations

import json
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

REPO = Path(__file__).resolve().parents[1]
RES = REPO / "hill395_experiment_results"
FIG = RES / "figures"
OUT = RES / "CITA-Net_백마고지_발표자료.pptx"

KFONT = "Malgun Gothic"      # Korean body
MONO = "Consolas"            # formulas
NAVY = RGBColor(0x1F, 0x2D, 0x3D)
BLUE = RGBColor(0x2C, 0x6F, 0xBB)
RED = RGBColor(0xC0, 0x39, 0x2B)
GREY = RGBColor(0x55, 0x5F, 0x6B)
LIGHT = RGBColor(0xF2, 0xF5, 0xF8)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

EMU_W, EMU_H = Inches(13.333), Inches(7.5)
prs = Presentation()
prs.slide_width = EMU_W
prs.slide_height = EMU_H
BLANK = prs.slide_layouts[6]

M = json.loads((RES / "metrics_summary.json").read_text(encoding="utf-8"))
ST = json.loads((RES / "dataset_stats.json").read_text(encoding="utf-8"))


def _set(run, size, color=NAVY, bold=False, font=KFONT, italic=False):
    run.font.size = Pt(size); run.font.bold = bold; run.font.italic = italic
    run.font.color.rgb = color; run.font.name = font


def slide():
    return prs.slides.add_slide(BLANK)


def bg(s, color):
    s.background.fill.solid(); s.background.fill.fore_color.rgb = color


def header(s, title, num=None):
    bar = s.shapes.add_shape(1, 0, 0, EMU_W, Inches(1.05))
    bar.fill.solid(); bar.fill.fore_color.rgb = NAVY; bar.line.fill.background()
    tf = bar.text_frame; tf.margin_left = Inches(0.45); tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]; r = p.add_run(); r.text = title
    _set(r, 26, WHITE, bold=True)
    if num is not None:
        tb = s.shapes.add_textbox(EMU_W - Inches(1.0), Inches(0.32), Inches(0.8), Inches(0.5))
        rr = tb.text_frame.paragraphs[0].add_run(); rr.text = str(num)
        _set(rr, 14, RGBColor(0xAA, 0xBB, 0xCC))


def textbox(s, x, y, w, h):
    tb = s.shapes.add_textbox(x, y, w, h); tb.text_frame.word_wrap = True
    return tb.text_frame


def bullets(tf, items, base=15):
    """items: list of (text, level, opts) ; opts dict optional."""
    first = True
    for it in items:
        text, lvl = it[0], it[1]
        opts = it[2] if len(it) > 2 else {}
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.level = lvl
        p.space_after = Pt(opts.get("after", 4)); p.space_before = Pt(opts.get("before", 2))
        prefix = "" if opts.get("nobull") else ("•  " if lvl == 0 else "–  ")
        r = p.add_run(); r.text = prefix + text
        _set(r, opts.get("size", base - lvl), opts.get("color", NAVY),
             bold=opts.get("bold", False), font=opts.get("font", KFONT))


def formula_box(s, x, y, w, lines, size=16, h=None):
    h = h or Inches(0.5 + 0.34 * len(lines))
    box = s.shapes.add_shape(1, x, y, w, h)
    box.fill.solid(); box.fill.fore_color.rgb = LIGHT
    box.line.color.rgb = BLUE; box.line.width = Pt(1)
    tf = box.text_frame; tf.word_wrap = True
    tf.margin_left = Inches(0.2); tf.margin_top = Inches(0.12)
    for i, ln in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(3)
        r = p.add_run(); r.text = ln
        _set(r, size, NAVY, font=MONO, bold=False)
    return box


def image_slide(title, num, img, caption, side_bullets=None):
    s = slide(); bg(s, WHITE); header(s, title, num)
    if side_bullets:
        im = s.shapes.add_picture(str(img), Inches(0.4), Inches(1.3), height=Inches(5.3))
        # clamp width
        if im.width > Inches(8.6):
            im.width = Inches(8.6); im.height = Inches(8.6 * im.height / im.width) if im.width else im.height
        tf = textbox(s, Inches(9.2), Inches(1.4), Inches(3.9), Inches(5.4))
        bullets(tf, side_bullets, base=14)
    else:
        s.shapes.add_picture(str(img), Inches(1.4), Inches(1.3), width=Inches(10.5))
        tf = textbox(s, Inches(0.6), Inches(6.7), Inches(12), Inches(0.6))
        r = tf.paragraphs[0].add_run(); r.text = caption; _set(r, 13, GREY, italic=True)
    return s


# ============================================================ 1. TITLE
s = slide(); bg(s, NAVY)
band = s.shapes.add_shape(1, 0, Inches(2.6), EMU_W, Inches(2.3))
band.fill.solid(); band.fill.fore_color.rgb = RGBColor(0x16, 0x22, 0x30); band.line.fill.background()
tf = textbox(s, Inches(0.8), Inches(2.7), Inches(11.7), Inches(2.1))
p = tf.paragraphs[0]; r = p.add_run()
r.text = "CITA-Net 기반 전장 STKG 개체정합(Entity Alignment)"; _set(r, 34, WHITE, bold=True)
p2 = tf.add_paragraph(); r = p2.add_run()
r.text = "Constraint-aware Identity & Track Association"; _set(r, 20, RGBColor(0x9F, 0xC5, 0xE8))
p3 = tf.add_paragraph(); p3.space_before = Pt(14); r = p3.add_run()
r.text = "사례 연구: 백마고지(Hill 395) 전투 시공간 지식그래프 — 방법론 · 수식 · 실험결과"
_set(r, 17, WHITE)
tf2 = textbox(s, Inches(0.8), Inches(6.6), Inches(11.7), Inches(0.6))
r = tf2.paragraphs[0].add_run()
r.text = "데이터셋 1,005,101 트리플 · 75 섹터 · cross-KG 관측 정합 · CPU 학습"
_set(r, 14, RGBColor(0xAA, 0xBB, 0xCC))

# ============================================================ 2. 문제정의
s = slide(); bg(s, WHITE); header(s, "1. 문제 정의 · 동기", 2)
tf = textbox(s, Inches(0.55), Inches(1.3), Inches(12.3), Inches(5.8))
bullets(tf, [
 ("과제: Cross-KG 개체정합(Entity Alignment) — 서로 다른 정보망이 따로 관측한 객체가 동일 실체인지 판정", 0, {"bold": True, "size": 18}),
 ("전장에서는 여러 센서·정보출처가 같은 부대를 비동기·부분적으로 관측한다", 1),
 ("KG-A: 아군 직접관측(전방관측조 VISUAL_OBS / RADAR / 포병관측 ARTILLERY_OBS)", 2),
 ("KG-B: 원거리·우회(항공정찰 AERIAL / 감청 SIGINT / 청음 ACOUSTIC / 정찰·포로 HUMINT)", 2),
 ("두 정보망의 관측을 융합하려면 '같은 실체'를 정확히 묶고(merge), 한쪽만 본 것은 매칭불가(dangling)로 남겨야 한다", 1),
 ("난점 (백마고지 도메인 그대로)", 0, {"bold": True, "size": 18, "before": 10}),
 ("동일 유형 부대 대량 등장(보병 소부대 분해: 돌격조·침투조·증원대·잔존대) → hard negative", 1),
 ("순간 위치만으로 구별 불가(능선 교차궤적), 한국어 라벨 다양·오식별, 비동기 시계(clock skew), 위치오차(CEP)", 1),
 ("Destroyed 잔해·Emplaced 화력 등 상태 제약, 운동학적으로 불가능한 전이 존재", 1),
 ("목표: 의미·시공간·관계·운동학·출처 제약을 한 모델에 통합해 정밀하면서 견고한 정합", 0,
  {"bold": True, "size": 18, "before": 10, "color": RED}),
], base=16)

# ============================================================ 3. 과제 형식화
s = slide(); bg(s, WHITE); header(s, "2. 과제의 형식적 정의", 3)
tf = textbox(s, Inches(0.55), Inches(1.25), Inches(12.3), Inches(2.2))
bullets(tf, [
 ("STKG = (Entities, Observations, Events, Landmarks, Relations) — 관측 중심 시공간 지식그래프", 0, {"bold": True}),
 ("관측 o = ⟨obs_id, kg, source, type, type_conf, state, state_conf, time, (E,N), cep_m, rels, events⟩", 1, {"font": MONO, "size": 13}),
 ("두 KG의 지역 개체 e^A, e^B 가 동일 실체이면 정합쌍, 어느 쪽에도 짝이 없으면 dangling", 1),
], base=15)
formula_box(s, Inches(0.7), Inches(3.7), Inches(12.0), [
 "입력 :  관측 집합 O = O_A ∪ O_B  (두 정보망)",
 "출력 :  정체성 군집 {I_k},  각 I_k ⊆ O  (관측→정체성 many-to-one)",
 "        + 매칭불가 집합 D ⊆ O  (∅ 슬롯으로 흡수)",
 "평가 :  교차-KG 개체쌍 (e^A, e^B) 을 정답과 비교 → Precision / Recall / F1",
], size=16)
tf = textbox(s, Inches(0.7), Inches(6.2), Inches(12), Inches(1.0))
bullets(tf, [("핵심 설계: 진영(ROK/CCF/US)을 객체 타입에 내장 → 진영 간 오병합을 구조적으로 차단하고, "
              "동일 진영·동일 기능 부대만 hard negative로 남긴다", 0, {"color": BLUE})], base=14)

# ============================================================ 4. 데이터셋
tdist = ST["type_distribution"]; sdist = ST["state_distribution"]
s = slide(); bg(s, WHITE); header(s, "3. 데이터셋 — 백마고지 Large Suite", 4)
tf = textbox(s, Inches(0.55), Inches(1.25), Inches(7.0), Inches(5.9))
bullets(tf, [
 (f"규모: {ST['total_triples']:,} 트리플 / {sum(ST['sectors'].values())} 섹터 / {ST['total_observations']:,} 관측", 0, {"bold": True}),
 (f"분할: train {ST['sectors']['train']} · dev {ST['sectors']['dev']} · test {ST['sectors']['test']} (분리 시드, 누수 안전)", 1),
 (f"knobs: identities/sector={ST['knobs']['identities_per_sector']}, obs/track={ST['knobs']['obs_per_track']}, "
  f"dangling={ST['knobs']['dangling_ratio']}", 1, {"size": 13}),
 ("생성 방식 (시드→절차적 확장)", 0, {"bold": True, "before": 10}),
 ("백마고지 사료(entity dict 386·taxonomy·.trig·Oct 6–15 타임라인)를 시드로 사용", 1),
 ("섹터=하위교전(국면×하위지역): ROK 방어 + CCF 공격 + 전차·포병·박격포·대공·공병·차량·조명", 1),
 ("소부대 분해→동일 type 대량(hard neg), 두 KG 이중관측+노이즈, dangling 20%", 1),
 ("결정론적 시드, stkg.nt(N-Triples)+observations.jsonl+gold_identities+dangling emit", 1),
 ("관측 스키마는 기존 large suite와 동일 → 파이프라인 드롭인", 0, {"color": BLUE, "before": 8}),
], base=15)
s.shapes.add_picture(str(FIG / "dataset_distributions.png"), Inches(7.7), Inches(2.6), width=Inches(5.4))
cap = textbox(s, Inches(7.7), Inches(6.0), Inches(5.4), Inches(0.5))
r = cap.paragraphs[0].add_run(); r.text = "타입/상태/소스별 관측 분포 (51,310 관측)"; _set(r, 12, GREY, italic=True)

# ============================================================ 5. 파이프라인 개요
s = slide(); bg(s, WHITE); header(s, "4. CITA-Net 파이프라인 개요", 5)
steps = [
 ("① Blocking", "그리드+시간\n후보 생성", BLUE),
 ("② Encoder", "5채널 다중모달\n+ 출처 게이트", BLUE),
 ("③ Relation GNN", "RGAT 2-hop\n관계맥락", BLUE),
 ("④ CTA", "제약인지\n전이 어텐션", RED),
 ("⑤ Decoder", "Sinkhorn\n슬롯 정합", BLUE),
 ("⑥ Heads/Loss", "pair·dangling\n·assign", GREY),
]
x = Inches(0.45); y = Inches(1.7); bw = Inches(2.0); gap = Inches(0.13)
for i, (t, d, c) in enumerate(steps):
    box = s.shapes.add_shape(1, x, y, bw, Inches(1.5))
    box.fill.solid(); box.fill.fore_color.rgb = c; box.line.fill.background()
    tf = box.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER; r = p.add_run(); r.text = t; _set(r, 15, WHITE, bold=True)
    p2 = tf.add_paragraph(); p2.alignment = PP_ALIGN.CENTER; r = p2.add_run(); r.text = d; _set(r, 11, WHITE)
    if i < len(steps) - 1:
        ar = s.shapes.add_shape(1, x + bw, y + Inches(0.6), gap, Inches(0.25))
        ar.fill.solid(); ar.fill.fore_color.rgb = GREY; ar.line.fill.background()
    x = x + bw + gap
tf = textbox(s, Inches(0.55), Inches(3.6), Inches(12.3), Inches(3.6))
bullets(tf, [
 ("관측 임베딩 h(o)을 인코더가 만들고, GNN이 관계맥락을 주입 → CTA가 시공간·운동학·상태·관계·출처 제약으로 전이확률 산출", 0),
 ("디코더가 K개 정체성 슬롯+∅ 슬롯에 관측을 Sinkhorn 정합 → 군집/궤적/전이 디코드", 0),
 ("학습은 pair·transition·trajectory·dangling·assign 5개 손실의 가중합으로 종단간(end-to-end)", 0),
 ("섹터 단위 스트리밍: 1M 트리플이지만 한 번에 한 섹터만 메모리 상주(peak ~16 MB)", 0, {"color": BLUE}),
], base=16)

# ============================================================ 6. Blocking
s = slide(); bg(s, WHITE); header(s, "5. 후보 생성 (Grid+Time Blocking)", 6)
tf = textbox(s, Inches(0.55), Inches(1.25), Inches(12.3), Inches(1.6))
bullets(tf, [
 ("O(M²) 전수 대신 (셀x, 셀y, 시간버킷) 공간해시로 시공간 이웃만 후보화", 0),
 ("관측별 reach 반경 = 운동학 도달거리 + 자기 CEP + 최악 파트너 CEP + 바닥값", 0),
 ("타입 호환(동일 카테고리 허용) + 전방 시간창(0 ≤ Δt ≤ dt_max) 필터", 0),
], base=15)
formula_box(s, Inches(0.7), Inches(3.1), Inches(12.0), [
 "reach_i = mult · v_max(type_i, state_i) · dt_max  +  cep_i + max_cep + r_floor + extra",
 "후보(i,j) 채택  ⇔  타입호환 ∧ (0 ≤ t_j − t_i ≤ dt_max) ∧ dist(i,j) ≤ reach",
 "설정:  dt_max = 180 s,  cell = 400 m,  r_floor = 80 m,  type_by_category = true",
], size=15)
tf = textbox(s, Inches(0.7), Inches(5.3), Inches(12), Inches(1.6))
bullets(tf, [
 ("결과: blocking recall ≥ 0.987, 후보 ~10× 축소, grid ⊆ brute-force(누락 없음)", 0, {"color": BLUE, "bold": True}),
 ("한국어 라벨 다양성 때문에 텍스트 게이트는 off(true match 손실 방지)", 0),
], base=15)

# ============================================================ 7. Encoder
s = slide(); bg(s, WHITE); header(s, "6. 관측 인코더 (Source-aware Multimodal)", 7)
tf = textbox(s, Inches(0.55), Inches(1.2), Inches(12.3), Inches(1.4))
bullets(tf, [
 ("관측 1개를 5개 채널로 임베딩 후, 출처 신뢰도 게이트 g(src)∈(0,1)⁵로 가중합", 0, {"bold": True}),
 ("채널: 텍스트(라벨) · 시공간(Fourier (x,y,t)) · 상태 · 타입(+conf) · 출처", 1),
], base=15)
formula_box(s, Inches(0.7), Inches(2.7), Inches(12.0), [
 "h(o) = LayerNorm( Σ_c  g_c(src) · f_c(o) ),   c ∈ {text, st, state, type, source}",
 "g(src) = σ( W · emb_gate(src) + logit(reliability(src)) )      # prior = sources.yaml",
 "f_st = MLP( Fourier(x, y, t) ),   Fourier = [sin(2^b·u), cos(2^b·u)]_b",
 "f_type = Linear( [emb_type(type) ;  type_confidence] )",
], size=15)
tf = textbox(s, Inches(0.7), Inches(5.2), Inches(12), Inches(1.8))
bullets(tf, [
 ("게이트 prior를 sources.yaml 신뢰도로 초기화 → RADAR는 위치↑/분류↓, HUMINT는 분류↑/위치↓ 등 출처 특성 반영", 0),
 ("d_model = 128, 텍스트 인코더=learned, 출처 게이트 on", 0, {"size": 14, "color": GREY}),
], base=15)

# ============================================================ 8. GNN
s = slide(); bg(s, WHITE); header(s, "7. 관계-맥락 그래프 인코더 (RGAT, 2-hop)", 8)
tf = textbox(s, Inches(0.55), Inches(1.2), Inches(12.3), Inches(1.5))
bullets(tf, [
 ("관측들이 이종 그래프를 형성: 동일트랙(시간연쇄) · 관계엣지(follows/near/supports/occupies…) · 공통이벤트", 0),
 ("다중관계 GAT 2-hop 메시지패싱으로 관계맥락을 임베딩에 흡수 (위치 모호 시 관계가 판별신호)", 0),
], base=15)
formula_box(s, Inches(0.7), Inches(2.8), Inches(12.0), [
 "msg(e) = W·h_src(e) + emb_rel(rel(e))                       # 관계별 메시지",
 "α_e = softmax_dst( LeakyReLU( aₛ·msg_h + a_d·(W·h_dst) ) )   # 엣지 어텐션(헤드별)",
 "h'_v = LayerNorm( h_v + ELU( Σ_{e→v} α_e · msg(e) ) )        # 잔차 + 정규화",
 "self-loop 관계 추가 · n_heads = 4 · layers = 2 · |REL| = 14(전투술어 3종 포함)",
], size=15)
tf = textbox(s, Inches(0.7), Inches(5.3), Inches(12), Inches(1.6))
bullets(tf, [
 ("전투 특화 술어 occupies/withdrawsFrom/reinforces를 별도 edge-type 임베딩으로 추가(featurize.REL_NAMES↔citanet.n_relations 연동)", 0),
 ("예: 'lead 전차를 follows 하는 UNKNOWN 보병'을 위치 모호에도 관계로 구별", 0, {"color": BLUE}),
], base=15)

# ============================================================ 9. Kinematics
s = slide(); bg(s, WHITE); header(s, "8. 운동학 타당성 (Kinematic Feasibility)", 9)
tf = textbox(s, Inches(0.55), Inches(1.25), Inches(12.3), Inches(1.2))
bullets(tf, [
 ("두 관측을 같은 객체로 이으려면 요구속도가 물리적 가능속도 이하여야 한다(CEP 위치오차 여유 포함)", 0, {"bold": True}),
], base=15)
formula_box(s, Inches(0.7), Inches(2.5), Inches(12.0), [
 "required_speed  =  dist(i,j) / Δt",
 "feasible_speed  =  max( v_max(type_i,state_i), v_max(type_j,state_j) ) + (cep_i + cep_j)/Δt",
 "feasible        ⇔  required_speed ≤ feasible_speed",
 "violation       =  max(0,  required_speed − feasible_speed)",
], size=16)
tf = textbox(s, Inches(0.7), Inches(4.9), Inches(12), Inches(2.0))
bullets(tf, [
 ("v_max(type,state)는 classes.yaml의 상태별 최대속도 — 예: 보병 2 m/s, 전차 12 m/s, 견인포 0, Emplaced/Holding은 0.5", 0),
 ("이 값이 CTA의 b_motion 항과 L_trajectory 손실, 그리고 blocking reach에 모두 사용됨", 0, {"color": BLUE}),
], base=15)

# ============================================================ 10. CTA (핵심)
s = slide(); bg(s, WHITE); header(s, "9. 제약인지 전이 어텐션 (CTA) — 핵심 수식", 10)
formula_box(s, Inches(0.55), Inches(1.25), Inches(12.3), [
 "score(i→j) = w₀·sim_sem + b_time + b_motion + b_state + b_rel + b_src",
 "p_transition(i→j) = σ( score(i→j) )",
], size=18, h=Inches(1.15))
formula_box(s, Inches(0.55), Inches(2.6), Inches(12.3), [
 "sim_sem  = cos( h_i , h_j )                              # 의미 유사도(인코더+GNN)",
 "b_time   = −BIG  if  t_j < t_i − ε   else  0             # 하드(역방향 금지, 비학습)",
 "b_motion = −α · softplus( (req_speed − feas_speed) / β ) # 운동학 위반 소프트 벌점",
 "b_state  = γ · log( C[s_i, s_j] + ε )                    # 상태 전이 호환성 행렬",
 "b_rel    = δ · jaccard( 이웃관계_i , 이웃관계_j )          # 관계 이웃 중첩 (M2+)",
 "b_src    = SrcBias[ src_i , src_j ]                      # 출처쌍 학습 바이어스",
], size=15)
tf = textbox(s, Inches(0.55), Inches(5.65), Inches(12.3), Inches(1.6))
bullets(tf, [
 ("학습 파라미터: w₀, α, β(=exp(logβ)), γ, δ, 출처바이어스 표 — 전부 데이터로 학습", 0),
 ("enabled_terms로 각 항 on/off → ablation이 설정 변경만으로 가능. 본 실험: [sem,time,motion,state,rel,src] 전부 사용", 0, {"color": BLUE}),
], base=14)

# ============================================================ 11. CTA 항 의미
s = slide(); bg(s, WHITE); header(s, "9-b. CTA 각 항의 역할", 11)
tf = textbox(s, Inches(0.55), Inches(1.3), Inches(12.3), Inches(5.8))
bullets(tf, [
 ("sim_sem (의미)", 0, {"bold": True, "color": BLUE}),
 ("인코더+GNN 임베딩의 코사인 — 같은 부대면 라벨·타입·관계맥락이 유사", 1),
 ("b_time (시간)", 0, {"bold": True, "color": BLUE}),
 ("후보는 전방시간만(t_i≤t_j); 역방향 전이는 −BIG로 사실상 금지(유일한 하드 제약)", 1),
 ("b_motion (운동학)", 0, {"bold": True, "color": BLUE}),
 ("요구속도가 가능속도를 넘으면 softplus 벌점 — Destroyed/Emplaced가 멀리 점프하는 오병합 억제", 1),
 ("b_state (상태)", 0, {"bold": True, "color": BLUE}),
 ("states.yaml의 비대칭 호환행렬 C — 예: Destroyed→이동 ≈0.01, Withdrawing→Occupying 낮음", 1),
 ("b_rel (관계)", 0, {"bold": True, "color": BLUE}),
 ("관계 이웃(술어+이벤트) 자카드 중첩 — 위치 모호 시 관계가 판별 신호", 1),
 ("b_src (출처)", 0, {"bold": True, "color": BLUE}),
 ("(출처_i, 출처_j) 쌍별 학습 바이어스 — 신뢰 높은 출처쌍의 정합을 우대", 1),
], base=15)

# ============================================================ 12. Decoder
s = slide(); bg(s, WHITE); header(s, "10. 정체성-궤적 디코더 (Sinkhorn 슬롯 정합)", 12)
tf = textbox(s, Inches(0.55), Inches(1.2), Inches(12.3), Inches(1.3))
bullets(tf, [
 ("K개 학습가능 정체성 슬롯 + 1개 ∅(null) 슬롯. 관측을 Sinkhorn 정규화 어텐션으로 슬롯에 소프트 배정", 0, {"bold": True}),
 ("CTA 전이확률로 이웃을 먼저 모아(컨텍스트) 강하게 연결된 관측이 같은 슬롯에 모이도록 유도", 0),
], base=15)
formula_box(s, Inches(0.7), Inches(2.75), Inches(12.0), [
 "h_ctx = h + σ(gate)·( Σ_j p_transition(i,j)·h_j  / Σ_j p_transition )   # CTA 컨텍스트",
 "Z = ( q(h_ctx) · k([slots ; null])ᵀ ) / √d                              # (M, K+1) 로짓",
 "A = Sinkhorn_row( Z / τ ) :  반복 [ 행 정규화 → 실슬롯 상대 열균형 ]      # 행합 = 1",
 "∅ 열은 균형서 제외 → dangling을 자유롭게 흡수.  디코드 = argmax_k A",
], size=15)
tf = textbox(s, Inches(0.7), Inches(5.2), Inches(12), Inches(1.7))
bullets(tf, [
 ("슬롯 내 관측을 kg별로 묶으면 1:1 정체성↔정체성 정합, 시간정렬하면 궤적/전이 산출", 0),
 ("설정: num_slots=48, sinkhorn_iters=20, τ=0.5, col_strength=0.1", 0, {"size": 14, "color": GREY}),
], base=15)

# ============================================================ 13. Heads + Loss
s = slide(); bg(s, WHITE); header(s, "11. 예측 헤드 & 손실 함수", 13)
tf = textbox(s, Inches(0.55), Inches(1.2), Inches(12.3), Inches(1.1))
bullets(tf, [
 ("PairHead: [h_i, h_j, |h_i−h_j|, h_i·h_j, CTA(6), (Δt,dist,text_cos)] → MLP → 교차-KG same 로짓", 0, {"size": 14}),
 ("DanglingHead: [h_o, (best/mean 후보로짓, 후보수, 최대유사도)] → MLP → 매칭불가 로짓", 0, {"size": 14}),
], base=14)
formula_box(s, Inches(0.7), Inches(2.5), Inches(12.0), [
 "L = λ_pair·L_pair + λ_trans·L_trans + λ_traj·L_traj + λ_dang·L_dang + λ_assign·L_assign",
 "L_pair   = BCE( pair_logits , same? )         # 교차-KG 후보쌍",
 "L_trans  = BCE( CTA.score , 동일정체성? )       # gold 할당에서 유도(대규모: 페어목록 無)",
 "L_traj   = mean( σ(pair_logit) · max(0, req−feas) ) / 10   # 소프트 운동학 일관성",
 "L_dang   = BCE( dangling_logits , dangling? )",
 "L_assign = CE( Z , 정체성↔슬롯 정렬된 타깃슬롯 )  # ∅ 슬롯=dangling",
], size=14)
tf = textbox(s, Inches(0.7), Inches(5.55), Inches(12), Inches(1.5))
bullets(tf, [
 ("가중치: λ_pair 2.0 · λ_trans 0.3 · λ_traj 0.5 · λ_dang 1.0 · λ_assign 2.0, pair_pos_weight 5.0", 0, {"size": 14}),
 ("L_assign은 Sinkhorn 할당(detach)에서 greedy 최대중첩으로 슬롯↔gold 정렬 후 관측별 교차엔트로피", 0, {"size": 14, "color": GREY}),
], base=14)

# ============================================================ 14. 학습 설정
s = slide(); bg(s, WHITE); header(s, "12. 학습 설정", 14)
tf = textbox(s, Inches(0.55), Inches(1.3), Inches(12.3), Inches(5.6))
bullets(tf, [
 ("설정 파일: configs/cita_full_hill395.yaml (data_root=battlefield_hill395_large)", 0, {"font": MONO, "size": 13}),
 ("모델: Full CITA-Net — 인코더(d=128, 출처게이트) + RGAT(2-hop,4-head) + CTA(6항) + Sinkhorn 디코더(48슬롯)", 0),
 ("최적화: Adam, lr 1e-3, grad-clip, 40 epoch, 매 epoch 전체 50 train 섹터 스트리밍", 0),
 ("blocking: dt_max 180s, cell 400m, r_floor 80m, type_by_category=true, text_gate=false", 0),
 ("학습 환경: 64-bit Python venv, CPU(torch 2.12 cpu), 약 45초/epoch", 0),
 ("학습 곡선: 총손실 17.63 → 1.63, best dev-F1 0.693 (모든 손실항 안정 수렴)", 0, {"color": BLUE, "bold": True, "before": 8}),
 ("스트리밍: 한 번에 한 섹터만 featurize → peak memory ~16 MB (1M 트리플 suite)", 0, {"color": BLUE}),
], base=16)

# ============================================================ 15. 평가지표 정의
s = slide(); bg(s, WHITE); header(s, "13. 평가 지표 정의", 15)
tf = textbox(s, Inches(0.55), Inches(1.3), Inches(12.3), Inches(5.7))
bullets(tf, [
 ("Precision / Recall / F1", 0, {"bold": True, "color": BLUE}),
 ("예측 교차-KG 개체쌍 (e^A,e^B)을 정답쌍과 비교 (TP/예측, TP/정답, 조화평균)", 1),
 ("Wrong-Merge Rate", 0, {"bold": True, "color": BLUE}),
 ("한 예측 정체성이 2개 이상의 서로 다른 gold 정체성을 포함한 비율(오병합)", 1),
 ("Fragmentation Rate", 0, {"bold": True, "color": BLUE}),
 ("한 gold 정체성이 여러 슬롯으로 쪼개지거나 회수되지 못한 비율(단편화/미회수)", 1),
 ("Trajectory Consistency / Impossible-Transition Rate", 0, {"bold": True, "color": BLUE}),
 ("디코드된 궤적의 전이 중 운동학적으로 불가능한 비율과 그 여집합", 1),
 ("Dangling Precision / Recall", 0, {"bold": True, "color": BLUE}),
 ("매칭불가(한쪽만 관측) 개체 탐지의 정밀도/재현율", 1),
], base=15)

# ============================================================ 16. 정량 결과 표 + 그림
dev, test = M["dev"], M["test"]
s = slide(); bg(s, WHITE); header(s, "14. 정량 결과 — dev / test 집계", 16)
rows = [("Precision", "precision"), ("Recall", "recall"), ("F1", "f1"),
        ("Wrong-merge", "wrong_merge_rate"), ("Fragmentation", "fragmentation_rate"),
        ("Traj-consistency", "trajectory_consistency_rate"),
        ("Dangling P", "dangling_precision"), ("Dangling R", "dangling_recall")]
tbl = s.shapes.add_table(len(rows) + 1, 3, Inches(0.55), Inches(1.3), Inches(5.6), Inches(5.2)).table
tbl.columns[0].width = Inches(2.8); tbl.columns[1].width = Inches(1.4); tbl.columns[2].width = Inches(1.4)
hdr = ["지표", "dev", "test"]
for c, t in enumerate(hdr):
    cell = tbl.cell(0, c); cell.text = t
    cell.fill.solid(); cell.fill.fore_color.rgb = NAVY
    p = cell.text_frame.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    _set(p.runs[0], 13, WHITE, bold=True)
for i, (lab, key) in enumerate(rows, start=1):
    vals = [lab, f"{dev[key]:.3f}", f"{test[key]:.3f}"]
    for c, v in enumerate(vals):
        cell = tbl.cell(i, c); cell.text = v
        p = cell.text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.LEFT if c == 0 else PP_ALIGN.CENTER
        _set(p.runs[0], 12, NAVY, bold=(key == "f1"))
        if i % 2 == 0:
            cell.fill.solid(); cell.fill.fore_color.rgb = LIGHT
s.shapes.add_picture(str(FIG / "dev_test_metrics.png"), Inches(6.4), Inches(1.5), width=Inches(6.6))
tf = textbox(s, Inches(6.4), Inches(5.6), Inches(6.6), Inches(1.5))
bullets(tf, [
 (f"test F1 {test['f1']:.3f} · Precision {test['precision']:.3f} · Recall {test['recall']:.3f}", 0, {"bold": True, "color": RED, "size": 15}),
 ("고정밀(P≈0.9)·중회수(R≈0.56): 동일 type hard negative와 노이즈를 의도적으로 넣은 결과(정상)", 0, {"size": 13}),
], base=13)

# ============================================================ 17. 학습곡선 + baseline
image_slide("15. 학습 곡선 & 기존 large suite 대비", 17, FIG / "training_curve.png",
            "", side_bullets=[
 ("학습 수렴", 0, {"bold": True, "color": BLUE}),
 ("총손실 17.63 → 1.63", 1),
 ("best dev-F1 → 0.693", 1),
 ("모든 손실항 안정 하강", 1),
 ("기존 large suite 대비", 0, {"bold": True, "color": RED, "before": 10}),
 ("baseline(README): F1~0.52 / P~0.80 / R~0.40", 1),
 ("백마고지 test: F1 0.686 / P 0.895 / R 0.561", 1),
 ("진영 내장 타입 + 충분한 epoch로 상회", 1),
 ("(baseline_compare.png 참조)", 1, {"size": 12, "color": GREY}),
])

# ============================================================ 18. baseline + 분포
image_slide("15-b. P/R/F1 비교 (baseline vs 백마고지)", 18, FIG / "baseline_compare.png", "",
            side_bullets=[
 ("회색 = 기존 large suite 참고치", 0),
 ("빨강 = 백마고지 test", 0),
 ("3개 지표 모두 상회", 0, {"color": RED, "bold": True}),
 ("특히 Precision 0.80→0.90", 1),
 ("Recall 0.40→0.56", 1),
])

# ============================================================ 19. 섹터별
image_slide("16. 섹터별 정밀도/재현율/F1", 19, FIG / "per_sector_metrics.png",
            "25개 dev+test 섹터별 성능 — 막대=F1, 선=precision/recall (점선 좌측 dev, 우측 test)")

# ============================================================ 20. 정성 결과
s = slide(); bg(s, WHITE); header(s, "17. 정성 결과 — 사례 분석", 20)
tf = textbox(s, Inches(0.55), Inches(1.25), Inches(7.4), Inches(5.9))
bullets(tf, [
 ("✅ 올바른 병합 (예)", 0, {"bold": True, "color": BLUE}),
 ("[KG-A] 중공군 침투조(CCF_Infantry, VISUAL/RADAR) ↔ [KG-B] 중공군 잔존대(CCF_Infantry, AERIAL/SIGINT) — 동일 실체 정합", 1, {"size": 12}),
 ("[KG-A] quad-50 대공포 ↔ [KG-B] quad-.50 — 확신도 1.0", 1, {"size": 12}),
 ("❌ 잘못된 병합 (예)", 0, {"bold": True, "color": RED, "before": 8}),
 ("동일 진영·동일 type 보병 소부대 간 오병합이 다수 (hard negative)", 1, {"size": 12}),
 ("단, 오병합은 확신도가 낮음(0.04~0.4) → 모델이 불확실해함", 1, {"size": 12}),
 ("✂️ 단편화 / 🚫 dangling", 0, {"bold": True, "color": GREY, "before": 8}),
 ("불확실하면 합치지 않고 abstain → 정밀도 우선, 재현율 손실", 1, {"size": 12}),
 ("한쪽 KG만 본 부대를 dangling으로 다수 정탐", 1, {"size": 12}),
 ("결론: hard negative가 실제로 어려운 케이스로 작동, 모델은 보수적", 0, {"color": BLUE, "before": 8, "size": 13}),
], base=14)
s.shapes.add_picture(str(FIG / "type_confusion.png"), Inches(8.1), Inches(1.6), width=Inches(5.0))
cap = textbox(s, Inches(8.1), Inches(5.9), Inches(5.0), Inches(0.6))
r = cap.paragraphs[0].add_run(); r.text = "오병합이 잦은 타입쌍 — 동일 진영 보병에 집중"; _set(r, 12, GREY, italic=True)

# ============================================================ 21. 결론
s = slide(); bg(s, WHITE); header(s, "18. 결론 & 향후 과제", 21)
tf = textbox(s, Inches(0.55), Inches(1.3), Inches(12.3), Inches(5.8))
bullets(tf, [
 ("기여", 0, {"bold": True, "size": 18, "color": BLUE}),
 ("백마고지 사료 기반 100만 트리플 전장 STKG를 CITA-Net 관측형으로 구축(드롭인) — 코드 변경 최소", 1),
 ("의미·시공간·운동학·상태·관계·출처 제약을 CTA 한 모듈에 통합, Sinkhorn 슬롯 디코더로 군집+dangling 동시 처리", 1),
 ("test F1 0.686 / P 0.895 — 기존 large suite 대비 향상, 고정밀·견고", 1),
 ("관찰", 0, {"bold": True, "size": 18, "color": BLUE, "before": 8}),
 ("오류는 동일 진영·동일 기능 소부대(hard negative)에 집중, 모델은 불확실 시 보수적 abstain", 1),
 ("향후", 0, {"bold": True, "size": 18, "color": RED, "before": 8}),
 ("재현율 향상(λ·블로킹·텍스트 인코더 튜닝), 원본 reified 이벤트-KG 충실 재현", 1),
 ("GeoSPARQL·OWL-Time·SKOS 표현으로 KG 병합/추론(누가 crest를 통제? 등) 실험 확장", 1),
 ("실제 좌표 기반 시각화·지도 오버레이로 작전 해석 지원", 1),
], base=15)

prs.save(str(OUT))
print("saved ->", OUT)
print("slides:", len(prs.slides._sldIdLst))

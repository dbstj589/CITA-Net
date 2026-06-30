# -*- coding: utf-8 -*-
"""Build the CITA-Net Hill 395 presentation (paper-flow, ~10 min, visual-heavy).

Run with the project venv:
    .venv/Scripts/python.exe scripts/build_ppt.py
"""
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

ROOT = Path(__file__).resolve().parent.parent
FIG = ROOT / "hill395_experiment_results" / "figures"
OUT = ROOT / "hill395_experiment_results" / "CITA-Net_백마고지_발표자료_v2.pptx"

# ----------------------------------------------------------------------------
# palette
NAVY   = RGBColor(0x1B, 0x2A, 0x4A)
BLUE   = RGBColor(0x2E, 0x5E, 0x8C)
STEEL  = RGBColor(0x4A, 0x7B, 0xB0)
SKY    = RGBColor(0xDCE, 0xE8, 0xF4) if False else RGBColor(0xDC, 0xE8, 0xF4)
ACCENT = RGBColor(0xE0, 0x7A, 0x28)   # orange
GREEN  = RGBColor(0x2E, 0x8B, 0x57)
RED    = RGBColor(0xC0, 0x39, 0x2B)
PURPLE = RGBColor(0x7D, 0x46, 0x9C)
GRAY   = RGBColor(0x5A, 0x64, 0x70)
LGRAY  = RGBColor(0xEC, 0xF0, 0xF4)
MGRAY  = RGBColor(0xC9, 0xD2, 0xDB)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
DARK   = RGBColor(0x22, 0x2A, 0x33)
KGA    = BLUE       # 아군 직접관측
KGB    = PURPLE     # 항공·감청

FONT = "맑은 고딕"
EMU_IN = 914400

# slide geometry (16:9)
SW, SH = 13.333, 7.5


def _set_font(run, size=14, bold=False, color=DARK, italic=False, name=FONT):
    f = run.font
    f.size = Pt(size)
    f.bold = bold
    f.italic = italic
    f.name = name
    f.color.rgb = color
    # ensure east-asian font
    rPr = run._r.get_or_add_rPr()
    ea = rPr.find(qn('a:ea'))
    if ea is None:
        ea = rPr.makeelement(qn('a:ea'), {})
        rPr.append(ea)
    ea.set('typeface', name)


def no_shadow(shape):
    try:
        shape.shadow.inherit = False
    except Exception:
        pass


def bg(slide, color=WHITE):
    r = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(SW), Inches(SH))
    r.fill.solid(); r.fill.fore_color.rgb = color
    r.line.fill.background(); no_shadow(r)
    slide.shapes._spTree.remove(r._element)
    slide.shapes._spTree.insert(2, r._element)
    return r


def rect(slide, x, y, w, h, fill=None, line=None, line_w=1.0,
         shape=MSO_SHAPE.RECTANGLE, radius=None):
    s = slide.shapes.add_shape(shape, Inches(x), Inches(y), Inches(w), Inches(h))
    if fill is None:
        s.fill.background()
    else:
        s.fill.solid(); s.fill.fore_color.rgb = fill
    if line is None:
        s.line.fill.background()
    else:
        s.line.color.rgb = line; s.line.width = Pt(line_w)
    no_shadow(s)
    if radius is not None and shape == MSO_SHAPE.ROUNDED_RECTANGLE:
        try:
            s.adjustments[0] = radius
        except Exception:
            pass
    return s


def txt(slide, x, y, w, h, lines, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
        wrap=True):
    """lines: list of dicts or list of [ (text,size,bold,color,align?) ]
    Each line is a dict: {t, s, b, c, a, sa(space_after pt), it}
    """
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = wrap
    tf.vertical_anchor = anchor
    for m in ('left', 'right', 'top', 'bottom'):
        setattr(tf, f'margin_{m}', Pt(2))
    first = True
    for ln in lines:
        if isinstance(ln, str):
            ln = {'t': ln}
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.alignment = ln.get('a', align)
        if ln.get('sa') is not None:
            p.space_after = Pt(ln['sa'])
        else:
            p.space_after = Pt(2)
        if ln.get('sb') is not None:
            p.space_before = Pt(ln['sb'])
        # support inline runs via 'runs'
        if 'runs' in ln:
            for r in ln['runs']:
                run = p.add_run(); run.text = r['t']
                _set_font(run, r.get('s', 14), r.get('b', False),
                          r.get('c', DARK), r.get('it', False))
        else:
            run = p.add_run(); run.text = ln['t']
            _set_font(run, ln.get('s', 14), ln.get('b', False),
                      ln.get('c', DARK), ln.get('it', False))
    return tb


def boxtext(shape, lines, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE):
    tf = shape.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    for m in ('left', 'right', 'top', 'bottom'):
        setattr(tf, f'margin_{m}', Pt(3))
    first = True
    for ln in lines:
        if isinstance(ln, str):
            ln = {'t': ln}
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.alignment = ln.get('a', align)
        p.space_after = Pt(ln.get('sa', 1))
        if 'runs' in ln:
            for r in ln['runs']:
                run = p.add_run(); run.text = r['t']
                _set_font(run, r.get('s', 13), r.get('b', False),
                          r.get('c', WHITE), r.get('it', False))
        else:
            run = p.add_run(); run.text = ln['t']
            _set_font(run, ln.get('s', 13), ln.get('b', False),
                      ln.get('c', WHITE), ln.get('it', False))
    return shape


def arrow(slide, x, y, w, h, color=STEEL, shape=MSO_SHAPE.RIGHT_ARROW):
    s = rect(slide, x, y, w, h, fill=color, shape=shape)
    return s


def chip(slide, x, y, w, text, fill, fg=WHITE, size=11, h=0.32, bold=True,
         radius=0.5):
    s = rect(slide, x, y, w, h, fill=fill, shape=MSO_SHAPE.ROUNDED_RECTANGLE,
             radius=radius)
    boxtext(s, [{'t': text, 's': size, 'b': bold, 'c': fg}])
    return s


# ----------------------------------------------------------------------------
prs = Presentation()
prs.slide_width = Inches(SW)
prs.slide_height = Inches(SH)
BLANK = prs.slide_layouts[6]


def new():
    s = prs.slides.add_slide(BLANK)
    bg(s, WHITE)
    return s


def header(slide, kicker, title, num):
    # left accent bar
    rect(slide, 0, 0, 0.18, SH, fill=ACCENT)
    rect(slide, 0.18, 0, SW - 0.18, 1.12, fill=NAVY)
    txt(slide, 0.55, 0.16, 9.5, 0.32,
        [{'t': kicker, 's': 12, 'b': True, 'c': RGBColor(0xF2, 0xC2, 0x8B)}])
    txt(slide, 0.52, 0.42, 11.5, 0.62,
        [{'t': title, 's': 25, 'b': True, 'c': WHITE}])
    # slide number badge
    b = rect(slide, SW - 1.0, 0.30, 0.55, 0.52, fill=ACCENT,
             shape=MSO_SHAPE.OVAL)
    boxtext(b, [{'t': str(num), 's': 16, 'b': True, 'c': WHITE}])
    # footer
    txt(slide, 0.55, SH - 0.38, 8, 0.3,
        [{'t': 'CITA-Net · 백마고지(Hill 395) STKG 개체정합', 's': 9, 'c': GRAY}])
    txt(slide, SW - 2.2, SH - 0.38, 1.7, 0.3,
        [{'t': f'{num} / 16', 's': 9, 'c': GRAY, 'a': PP_ALIGN.RIGHT}])


def section_tag(slide, x, y, label, color=ACCENT):
    chip(slide, x, y, 1.5, label, color, size=11)


# helper: small icon circle with glyph
def icon(slide, x, y, d, glyph, fill, fg=WHITE, size=18):
    c = rect(slide, x, y, d, d, fill=fill, shape=MSO_SHAPE.OVAL)
    boxtext(c, [{'t': glyph, 's': size, 'b': True, 'c': fg}])
    return c


# motif: two KGs observing same battlefield
def two_kg_motif(slide, x, y, w=4.0, label_a='KG-A 아군 직접관측',
                 label_b='KG-B 항공·감청', a_text='국군 보병 대대',
                 b_text='미상 보병', q=True, scale=1.0):
    bw = w
    bh = 0.62 * scale
    a = rect(slide, x, y, bw, bh, fill=KGA, shape=MSO_SHAPE.ROUNDED_RECTANGLE,
             radius=0.18)
    boxtext(a, [{'t': label_a, 's': 11, 'b': True, 'c': WHITE},
                {'t': f'“{a_text}”', 's': 12, 'b': False, 'c': RGBColor(0xDC,0xE8,0xF4)}])
    b = rect(slide, x, y + bh + 0.55 * scale, bw, bh, fill=KGB,
             shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.18)
    boxtext(b, [{'t': label_b, 's': 11, 'b': True, 'c': WHITE},
                {'t': f'“{b_text}”', 's': 12, 'b': False, 'c': RGBColor(0xEAD if False else 0xEA,0xDD,0xF2)}])
    if q:
        qm = rect(slide, x + bw + 0.18, y + bh * 0.55, 0.7, 0.7, fill=ACCENT,
                  shape=MSO_SHAPE.OVAL)
        boxtext(qm, [{'t': '=?', 's': 18, 'b': True, 'c': WHITE}])
    return a, b


# ============================================================ SLIDE 1: TITLE
s = new()
rect(s, 0, 0, SW, SH, fill=NAVY)
rect(s, 0, 0, SW, 0.22, fill=ACCENT)
rect(s, 0, SH - 0.22, SW, 0.22, fill=ACCENT)
# decorative faint grid blocks (battlefield sectors)
for i in range(6):
    rect(s, 8.6 + (i % 3) * 1.45, 1.1 + (i // 3) * 1.45, 1.3, 1.3,
         fill=None, line=RGBColor(0x33, 0x46, 0x66), line_w=1.0)
txt(s, 0.9, 1.7, 8.5, 0.5,
    [{'t': '국방 시공간 지식그래프 · 개체정합', 's': 16, 'b': True,
      'c': RGBColor(0xF2, 0xC2, 0x8B)}])
txt(s, 0.85, 2.25, 9.5, 1.8,
    [{'t': 'CITA-Net', 's': 54, 'b': True, 'c': WHITE, 'sa': 4},
     {'t': '제약-인지 전이 어텐션 기반', 's': 23, 'b': True, 'c': WHITE, 'sa': 2},
     {'t': '다중 정보망 전장 개체정합 신경망', 's': 23, 'b': True, 'c': WHITE}])
txt(s, 0.9, 4.6, 9, 0.9,
    [{'t': '백마고지(Hill 395) 전투 시나리오 기반 실험', 's': 15,
      'c': RGBColor(0xCF, 0xDA, 0xE8), 'sa': 3},
     {'t': 'Constraint-aware Identity-Trajectory Alignment Network', 's': 13,
      'it': True, 'c': RGBColor(0x9F, 0xB2, 0xC8)}])
rect(s, 0.92, 5.95, 3.2, 0.04, fill=ACCENT)
txt(s, 0.9, 6.15, 9, 0.5,
    [{'t': '데이터셋 1,005,101 트리플 · 75 섹터 · Full CITA-Net (M3)', 's': 12,
      'c': RGBColor(0xAF, 0xC0, 0xD4)}])

# ============================================================ SLIDE 2: AGENDA
s = new()
header(s, 'CONTENTS', '발표 목차 — 논문 구성', 2)
items = [
    ('1', '개요', 'Overview', '문제 정의: 두 정보망이 본 같은 전장', NAVY),
    ('2', '관련 연구', 'Related Work', '엔티티 정합 최신 연구와 한계', BLUE),
    ('3', '방법론', 'Methodology', '6개 모듈 × 전장 예시로 직관 설명', ACCENT),
    ('4', '실험 결과', 'Experiments', '정량·정성 분석 (백마고지)', GREEN),
    ('5', '결론', 'Conclusion', '기여점 · 한계 · 향후 연구', PURPLE),
]
y = 1.5
for n, ko, en, desc, col in items:
    c = rect(s, 0.8, y, 0.75, 0.75, fill=col, shape=MSO_SHAPE.OVAL)
    boxtext(c, [{'t': n, 's': 26, 'b': True, 'c': WHITE}])
    bar = rect(s, 1.75, y + 0.04, 10.7, 0.67, fill=LGRAY,
               shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.12)
    txt(s, 2.0, y + 0.06, 4.0, 0.65,
        [{'runs': [{'t': ko + '  ', 's': 18, 'b': True, 'c': NAVY},
                   {'t': en, 's': 11, 'b': False, 'c': GRAY}]}],
        anchor=MSO_ANCHOR.MIDDLE)
    txt(s, 6.2, y + 0.06, 6.0, 0.65, [{'t': desc, 's': 13, 'c': DARK}],
        anchor=MSO_ANCHOR.MIDDLE)
    y += 1.02

# ============================================================ SLIDE 3: OVERVIEW problem
s = new()
header(s, '1. 개요  |  OVERVIEW', '문제 정의 — 같은 전장, 서로 다른 두 개의 보고', 3)
txt(s, 0.6, 1.3, 12, 0.5,
    [{'runs': [
        {'t': '백마고지 전투', 's': 15, 'b': True, 'c': ACCENT},
        {'t': '에서 ', 's': 15, 'c': DARK},
        {'t': '두 정보망', 's': 15, 'b': True, 'c': NAVY},
        {'t': '이 같은 부대를 ', 's': 15, 'c': DARK},
        {'t': '다른 이름·좌표·시각', 's': 15, 'b': True, 'c': RED},
        {'t': '으로 보고한다.', 's': 15, 'c': DARK}]}])

# battlefield box center
field = rect(s, 4.55, 2.1, 4.25, 3.4, fill=RGBColor(0xEE,0xF2,0xE9),
             line=GREEN, line_w=1.5, shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.06)
txt(s, 4.55, 2.2, 4.25, 0.4, [{'t': '⛰  백마고지 (Hill 395)', 's': 13, 'b': True,
    'c': GREEN, 'a': PP_ALIGN.CENTER}])
# unit token in field
for (ux, uy, t) in [(5.6, 3.1, '보병'), (6.9, 3.0, '전차'), (6.0, 4.3, '포병')]:
    u = rect(s, ux, uy, 0.95, 0.55, fill=GRAY, shape=MSO_SHAPE.ROUNDED_RECTANGLE,
             radius=0.3)
    boxtext(u, [{'t': t, 's': 11, 'b': True, 'c': WHITE}])
txt(s, 4.55, 4.95, 4.25, 0.45,
    [{'t': '실제로는 하나의 전장 = 같은 부대들', 's': 11, 'it': True,
      'c': GRAY, 'a': PP_ALIGN.CENTER}])

# KG-A left
a = rect(s, 0.7, 2.35, 3.4, 1.25, fill=KGA, shape=MSO_SHAPE.ROUNDED_RECTANGLE,
         radius=0.08)
boxtext(a, [{'t': 'KG-A  아군 직접관측', 's': 13, 'b': True, 'c': WHITE},
            {'t': '포병관측 · 시각 · 레이더', 's': 10, 'c': RGBColor(0xCF,0xDD,0xEC)},
            {'t': '“국군 보병 대대 @좌표X”', 's': 11, 'b': True, 'c': WHITE}])
# KG-B right
b = rect(s, 9.25, 2.35, 3.4, 1.25, fill=KGB, shape=MSO_SHAPE.ROUNDED_RECTANGLE,
         radius=0.08)
boxtext(b, [{'t': 'KG-B  항공·감청', 's': 13, 'b': True, 'c': WHITE},
            {'t': '항공정찰 · SIGINT · 음향 · HUMINT', 's': 10, 'c': RGBColor(0xEA,0xDD,0xF2)},
            {'t': '“미상 보병 @좌표Y”', 's': 11, 'b': True, 'c': WHITE}])
arrow(s, 4.12, 2.8, 0.4, 0.35, color=KGA)
arrow(s, 8.8, 2.8, 0.4, 0.35, color=KGB, shape=MSO_SHAPE.LEFT_ARROW)

# core question banner
qb = rect(s, 1.7, 5.75, 9.9, 1.0, fill=NAVY, shape=MSO_SHAPE.ROUNDED_RECTANGLE,
          radius=0.1)
boxtext(qb, [{'runs': [
    {'t': '핵심 과제:  ', 's': 15, 'b': True, 'c': RGBColor(0xF2,0xC2,0x8B)},
    {'t': 'A의 “국군 보병 대대”와 B의 “미상 보병”은 ', 's': 15, 'c': WHITE},
    {'t': '같은 부대인가, 다른 부대인가?', 's': 15, 'b': True, 'c': WHITE}]},
    {'t': '→  이 자동 판정이 곧 “개체정합(Entity Alignment)”', 's': 12,
     'c': RGBColor(0xCF,0xDD,0xEC)}])

# ============================================================ SLIDE 4: why hard
s = new()
header(s, '1. 개요  |  OVERVIEW', '왜 어려운가 — 판정의 3가지 결과와 3대 난점', 4)
# three outcomes
txt(s, 0.6, 1.25, 12, 0.35, [{'t': '▣  판정은 세 갈래로 나뉜다', 's': 14, 'b': True, 'c': NAVY}])
outs = [
    ('✔ 병합 (Merge)', '같은 실체 → 하나로 통합', GREEN),
    ('✂ 단편화 (Fragment)', '같은데 확신 못 해 못 합침 (놓침)', ACCENT),
    ('⊘ 매칭불가 (Dangling)', '한쪽에만 잡힌 부대 → 짝 없음', RED),
]
x = 0.7
for t, d, col in outs:
    card = rect(s, x, 1.65, 3.9, 1.0, fill=WHITE, line=col, line_w=2.0,
                shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.08)
    txt(s, x + 0.18, 1.78, 3.6, 0.9,
        [{'t': t, 's': 14, 'b': True, 'c': col, 'sa': 3},
         {'t': d, 's': 11, 'c': DARK}])
    x += 4.1

txt(s, 0.6, 3.0, 12, 0.35, [{'t': '▣  기존 방법으로는 풀기 어려운 3대 난점', 's': 14, 'b': True, 'c': NAVY}])
hards = [
    ('🛰', '잡음 많은 다중 센서', '레이더·감청·인간정보가 제각각.\n출처마다 신뢰도가 다르다.', BLUE),
    ('🏃', '물리적 이동 제약', '부대는 순간이동 불가.\n시속을 넘는 “전이”는 불가능.', ACCENT),
    ('🔗', '관계 맥락 의존', '“미상 표적”은 단독으론 모름.\n‘무엇을 따라가는가’로만 식별.', PURPLE),
]
x = 0.7
for g, t, d, col in hards:
    card = rect(s, x, 3.45, 3.9, 2.7, fill=LGRAY,
                shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.06)
    ic = rect(s, x + 1.45, 3.65, 1.0, 1.0, fill=col, shape=MSO_SHAPE.OVAL)
    boxtext(ic, [{'t': g, 's': 30, 'c': WHITE}])
    txt(s, x + 0.2, 4.8, 3.5, 1.3,
        [{'t': t, 's': 14, 'b': True, 'c': col, 'a': PP_ALIGN.CENTER, 'sa': 4}] +
        [{'t': line, 's': 11, 'c': DARK, 'a': PP_ALIGN.CENTER}
         for line in d.split('\n')])
    x += 4.1
txt(s, 0.6, 6.35, 12, 0.5,
    [{'runs': [{'t': '⇒  CITA-Net은 이 3대 난점을 ', 's': 13, 'b': True, 'c': DARK},
               {'t': '출처게이트 · 운동학 제약(CTA) · 관계 GNN', 's': 13, 'b': True, 'c': ACCENT},
               {'t': '으로 정면 해결', 's': 13, 'b': True, 'c': DARK}]}])

# ============================================================ SLIDE 5: related work
s = new()
header(s, '2. 관련 연구  |  RELATED WORK', '엔티티 정합 연구의 3계열과 공통 한계', 5)
cols = [
    ('① 구조·임베딩 기반', 'MTransE, BootEA,\nGCN-Align, RDGCN,\nAliNet, Dual-AMN', BLUE,
     '두 그래프가 구조적으로\n닮았다고 가정.\n잡음·희소 그래프에 취약.'),
    ('② 멀티모달 정합', 'MMEA, EVA,\nMSNEA, MCLEA', PURPLE,
     '백과사전 KG(이미지·속성)용.\n시공간·운동학 추론 없음.\n출처 신뢰도 미반영.'),
    ('③ 시간·정합불가 처리', 'TEA-GNN, TREA,\nDangling-EA\n(Sun et al. 2021)', GREEN,
     '시간을 이산 타임스탬프로만.\n연속 기동(이동) 모델 없음.'),
]
x = 0.7
for t, methods, col, lim in cols:
    head = rect(s, x, 1.35, 3.9, 0.55, fill=col,
                shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.1)
    boxtext(head, [{'t': t, 's': 13, 'b': True, 'c': WHITE}])
    body = rect(s, x, 1.95, 3.9, 1.55, fill=WHITE, line=col, line_w=1.5,
                shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.05)
    txt(s, x + 0.15, 2.05, 3.6, 1.45,
        [{'t': m, 's': 12, 'b': True, 'c': DARK, 'a': PP_ALIGN.CENTER}
         for m in methods.split('\n')], anchor=MSO_ANCHOR.MIDDLE)
    limbox = rect(s, x, 3.62, 3.9, 1.5, fill=RGBColor(0xFB,0xEC,0xE6),
                  shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.05)
    txt(s, x + 0.15, 3.72, 3.6, 1.4,
        [{'t': '⚠ 한계', 's': 11, 'b': True, 'c': RED, 'a': PP_ALIGN.CENTER, 'sa': 3}] +
        [{'t': line, 's': 11, 'c': DARK, 'a': PP_ALIGN.CENTER}
         for line in lim.split('\n')])
    x += 4.1
# common gap banner
gap = rect(s, 0.7, 5.35, 11.9, 1.45, fill=NAVY,
           shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.06)
txt(s, 0.95, 5.5, 11.4, 1.3,
    [{'t': '🔻  공통 한계 (전장 STKG 관점)', 's': 13, 'b': True,
      'c': RGBColor(0xF2,0xC2,0x8B), 'sa': 4},
     {'runs': [
        {'t': '① 물리·운동학 제약 부재   ', 's': 12.5, 'c': WHITE},
        {'t': '② 출처 신뢰도 미반영   ', 's': 12.5, 'c': WHITE},
        {'t': '③ 정합불가(dangling) 미처리   ', 's': 12.5, 'c': WHITE},
        {'t': '④ 설명불가(블랙박스)', 's': 12.5, 'c': WHITE}], 'sa': 4},
     {'runs': [
        {'t': '⇒  CITA-Net: ', 's': 12.5, 'b': True, 'c': ACCENT},
        {'t': '멀티모달 융합 + 관계 GNN + 운동학 제약 + dangling 헤드를 ', 's': 12.5, 'c': WHITE},
        {'t': '설명가능한 가산식', 's': 12.5, 'b': True, 'c': WHITE},
        {'t': '으로 통합', 's': 12.5, 'c': WHITE}]}])

# ============================================================ SLIDE 6: methodology pipeline
s = new()
header(s, '3. 방법론  |  METHODOLOGY', '전체 파이프라인 — 개별→관계→제약→종합', 6)
txt(s, 0.6, 1.25, 12.2, 0.4,
    [{'t': '맥락을 한 겹씩 쌓아간다: 각 단계는 앞 단계의 출력(128차원 벡터)을 입력으로 받는다.',
      's': 13, 'c': DARK}])
stages = [
    ('①', '후보생성', 'blocking_grid\n불가능 쌍 가지치기', STEEL),
    ('②', '특징화', 'featurize\n관측·쌍·그래프 텐서', STEEL),
    ('③', '관측 인코더', '5채널→128\n출처 게이트', BLUE),
    ('④', '관계 GNN', 'RGAT ×2\n이웃으로 식별', BLUE),
    ('⑤', 'CTA', '6항 제약 점수\n(모델의 두뇌)', ACCENT),
    ('⑥', '헤드·디코더', 'Pair/Dangling\nSinkhorn 슬롯', GREEN),
]
n = len(stages)
gap = 0.18
w = (12.4 - gap * (n - 1)) / n
x = 0.7
y = 1.95
for i, (num, t, d, col) in enumerate(stages):
    c = rect(s, x, y, w, 2.0, fill=col, shape=MSO_SHAPE.ROUNDED_RECTANGLE,
             radius=0.08)
    boxtext(c, [{'t': num, 's': 22, 'b': True, 'c': WHITE, 'sa': 3},
                {'t': t, 's': 13, 'b': True, 'c': WHITE, 'sa': 4}] +
               [{'t': line, 's': 10, 'c': RGBColor(0xEC,0xF1,0xF6)}
                for line in d.split('\n')])
    if i < n - 1:
        arrow(s, x + w - 0.02, y + 0.78, gap + 0.06, 0.42, color=MGRAY)
    x += w + gap
# bottom rationale strip
labels = ['가지치기', '수치화', '개별 이해', '관계 반영', '물리 검증', '최종 판정']
x = 0.7
for i, lb in enumerate(labels):
    chip(s, x, y + 2.15, w, lb, LGRAY, fg=NAVY, size=10, bold=True)
    x += w + gap
txt(s, 0.7, 4.7, 12, 1.6,
    [{'t': '왜 이 순서인가?', 's': 14, 'b': True, 'c': NAVY, 'sa': 5},
     {'runs': [{'t': '①  ', 's': 12.5, 'b': True, 'c': STEEL},
               {'t': 'N² 쌍 폭발을 막아 그럴듯한 후보만 남기고  →  ', 's': 12.5, 'c': DARK},
               {'t': '③④  ', 's': 12.5, 'b': True, 'c': BLUE},
               {'t': '관측을 의미·관계로 “이해”한 뒤  →  ', 's': 12.5, 'c': DARK},
               {'t': '⑤  ', 's': 12.5, 'b': True, 'c': ACCENT},
               {'t': '“의미는 맞지만 물리적으로 불가능”을 가려내고  →  ', 's': 12.5, 'c': DARK},
               {'t': '⑥  ', 's': 12.5, 'b': True, 'c': GREEN},
               {'t': '모든 맥락이 쌓인 뒤 최종 결정.', 's': 12.5, 'c': DARK}], 'sa': 5},
     {'runs': [{'t': '핵심: ', 's': 12.5, 'b': True, 'c': RED},
               {'t': '모든 벡터를 128차원으로 통일', 's': 12.5, 'b': True, 'c': DARK},
               {'t': '해야 채널을 더하고(③), 부대끼리 비교하고(⑤), 다음 모듈로 흘려보낼 수 있다(④⑥).',
                's': 12.5, 'c': DARK}]}])

# ============================================================ SLIDE 7: blocking + featurize
s = new()
header(s, '3. 방법론  |  ①②  후보생성 · 특징화', '계산 폭발을 막는 “격자 가지치기”', 7)
txt(s, 0.6, 1.25, 12, 0.4,
    [{'runs': [{'t': '문제: ', 's': 13, 'b': True, 'c': RED},
               {'t': '관측 N개면 가능한 쌍은 N²개(수만 개). 전부 비교하면 폭발.', 's': 13, 'c': DARK}]}])
# grid visual left
gx, gy, cell = 0.9, 2.0, 0.72
for r in range(4):
    for cc in range(4):
        rect(s, gx + cc*cell, gy + r*cell, cell, cell, fill=None,
             line=MGRAY, line_w=1.0)
# highlight one cell + neighbors
rect(s, gx + cell, gy + cell, cell, cell, fill=RGBColor(0xFB,0xE3,0xCF),
     line=ACCENT, line_w=2.0)
# tokens
import random
random.seed(3)
for (cx, cy, col) in [(gx+1.3, gy+1.25, KGA), (gx+1.55, gy+1.55, KGB),
                      (gx+0.3, gy+0.25, KGA), (gx+2.6, gy+2.7, KGB),
                      (gx+1.4, gy+2.6, KGA)]:
    rect(s, cx, cy, 0.16, 0.16, fill=col, shape=MSO_SHAPE.OVAL)
txt(s, gx, gy + 4*cell + 0.05, 4*cell, 0.6,
    [{'t': '400m 격자 · ±3분 윈도 안의 쌍만 후보', 's': 11, 'c': GRAY,
      'a': PP_ALIGN.CENTER}])
# right: rules + featurize output
txt(s, 4.4, 1.9, 8.3, 0.4, [{'t': '① 후보생성 (blocking_grid)', 's': 14, 'b': True, 'c': BLUE}])
for i, (rule, val) in enumerate([
        ('시간차 ≤ 180초(3분)', '너무 떨어지면 제외'),
        ('400m 격자 + 이웃칸', '먼 거리 제외'),
        ('도달 가능성(필요속도)', '음속으로 달려야 하면 제외')]):
    chip(s, 4.5, 2.4 + i*0.5, 3.6, rule, LGRAY, fg=NAVY, size=11, h=0.4)
    txt(s, 8.25, 2.4 + i*0.5, 4.3, 0.4, [{'t': '→ ' + val, 's': 11, 'c': DARK}],
        anchor=MSO_ANCHOR.MIDDLE)
txt(s, 4.4, 4.15, 8.3, 0.4, [{'t': '② 특징화 (featurize) — 모델이 먹을 텐서로', 's': 14, 'b': True, 'c': BLUE}])
fcards = [
    ('정규화', '좌표 −중심 ÷1000\n→ 0 근처 작은 값', STEEL),
    ('쌍 제약 사전계산', '시간차·거리·필요속도\n상태호환·관계겹침', STEEL),
    ('교차-KG 관계비교', 'ID 대신 술어로\n(follows/near…)', STEEL),
]
x = 4.5
for t, d, col in fcards:
    card = rect(s, x, 4.6, 2.65, 1.45, fill=WHITE, line=col, line_w=1.3,
                shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.06)
    txt(s, x+0.12, 4.7, 2.45, 1.3,
        [{'t': t, 's': 12, 'b': True, 'c': col, 'a': PP_ALIGN.CENTER, 'sa': 3}] +
        [{'t': l, 's': 10.5, 'c': DARK, 'a': PP_ALIGN.CENTER} for l in d.split('\n')])
    x += 2.78
txt(s, 0.6, 6.45, 12, 0.4,
    [{'runs': [{'t': '전장 예시:  ', 's': 12, 'b': True, 'c': ACCENT},
               {'t': '3분 전 5km 밖 관측은 보병이 닿을 수 없어 후보에서 자동 제외 → 정밀 분석은 소수 쌍만.',
                's': 12, 'c': DARK}]}])

# ============================================================ SLIDE 8: encoder
s = new()
header(s, '3. 방법론  |  ③  관측 인코더', '5채널을 128차원으로 융합 + 출처 신뢰도 게이트', 8)
txt(s, 0.6, 1.22, 12, 0.4,
    [{'t': '관측 하나에 섞인 5종류의 이질적 정보를 “공통 양식(128차원)”으로 번역해 가중합한다.',
      's': 13, 'c': DARK}])
chans = [
    ('① 텍스트', '“국군 보병 대대”', BLUE),
    ('② 시공간', '좌표(x,y)+시각 t', STEEL),
    ('③ 상태', 'Engaging(교전중)', STEEL),
    ('④ 타입', 'ROK_Infantry', STEEL),
    ('⑤ 출처', 'RADAR', PURPLE),
]
y = 1.85
for i, (t, ex, col) in enumerate(chans):
    c = rect(s, 0.7, y, 3.2, 0.62, fill=col, shape=MSO_SHAPE.ROUNDED_RECTANGLE,
             radius=0.12)
    boxtext(c, [{'runs': [{'t': t + '  ', 's': 12, 'b': True, 'c': WHITE},
                          {'t': ex, 's': 10.5, 'c': RGBColor(0xE8,0xEE,0xF4)}]}],
            align=PP_ALIGN.LEFT)
    arrow(s, 3.95, y + 0.12, 0.5, 0.38, color=MGRAY)
    # each → 128 box
    p = rect(s, 4.5, y, 1.5, 0.62, fill=WHITE, line=col, line_w=1.3,
             shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.12)
    boxtext(p, [{'t': '→ 128차원', 's': 11, 'b': True, 'c': col}])
    y += 0.72
# gate
gate = rect(s, 6.25, 1.95, 1.9, 3.0, fill=ACCENT,
            shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.06)
boxtext(gate, [{'t': '출처 게이트', 's': 13, 'b': True, 'c': WHITE, 'sa': 4},
               {'t': 'g(src) ∈ (0,1)', 's': 11, 'c': WHITE, 'sa': 6},
               {'t': '출처에 따라', 's': 10.5, 'c': WHITE},
               {'t': '믿을 채널을', 's': 10.5, 'c': WHITE},
               {'t': '자동 조절', 's': 10.5, 'c': WHITE}])
arrow(s, 8.2, 3.2, 0.5, 0.5, color=MGRAY)
# fused
fused = rect(s, 8.8, 2.5, 3.8, 1.9, fill=NAVY,
             shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.06)
boxtext(fused, [{'t': 'h(o) = LayerNorm(Σ g·f )', 's': 13, 'b': True, 'c': WHITE, 'sa': 5},
                {'t': '관측 1개 = 의미 벡터', 's': 12, 'c': RGBColor(0xCF,0xDD,0xEC), 'sa': 2},
                {'t': '[0.2, −1.3, 0.7, … ] (128)', 's': 11, 'it': True,
                 'c': RGBColor(0xF2,0xC2,0x8B)}])
# example callout
ex = rect(s, 6.25, 5.15, 6.35, 1.55, fill=RGBColor(0xFB,0xEC,0xE6),
          shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.05)
txt(s, 6.45, 5.27, 6.0, 1.4,
    [{'t': '🛰  전장 예시 — 출처마다 믿을 정보가 다르다', 's': 12, 'b': True, 'c': ACCENT, 'sa': 4},
     {'runs': [{'t': 'RADAR', 's': 11.5, 'b': True, 'c': PURPLE},
               {'t': ' → 좌표는 정확하나 부대명은 모름 (시공간↑, 텍스트↓)', 's': 11.5, 'c': DARK}], 'sa': 2},
     {'runs': [{'t': 'HUMINT', 's': 11.5, 'b': True, 'c': PURPLE},
               {'t': ' → 명칭은 알지만 위치 부정확 (텍스트↑, 시공간↓)', 's': 11.5, 'c': DARK}]}])
lnote = rect(s, 0.7, 5.5, 5.4, 1.2, fill=LGRAY,
             shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.06)
txt(s, 0.9, 5.62, 5.0, 1.0,
    [{'t': '💡 사전지식으로 똑똑하게 시작', 's': 12, 'b': True, 'c': NAVY, 'sa': 3},
     {'t': '게이트 초기값을 ontology의 출처 신뢰도로 설정한 뒤, 데이터로 미세조정.',
      's': 11.5, 'c': DARK}])

# ============================================================ SLIDE 9: GNN
s = new()
header(s, '3. 방법론  |  ④  관계-맥락 GNN', '혼자선 모를 정체를 “이웃 관계”로 밝힌다', 9)
txt(s, 0.6, 1.22, 12, 0.4,
    [{'t': 'GNN: 관측을 점, 관계를 선으로 본 그래프 위에서 메시지를 주고받아 각 벡터를 갱신한다.',
      's': 13, 'c': DARK}])
# graph visual
nodes = {
    'lead': (3.3, 2.6, '미군 전차\n(US_Tank)', BLUE),
    'unk':  (3.3, 4.4, '미상 표적\n(UNKNOWN)', GRAY),
    'inf':  (1.3, 3.5, '중공군 보병', PURPLE),
    'col':  (5.3, 3.5, '차량대', STEEL),
}
# edges
def line_between(a, b, color=MGRAY, label=None, w=2.0):
    x1, y1 = nodes[a][0]+0.55, nodes[a][1]+0.35
    x2, y2 = nodes[b][0]+0.55, nodes[b][1]+0.35
    conn = s.shapes.add_connector(2, Inches(x1), Inches(y1), Inches(x2), Inches(y2))
    conn.line.color.rgb = color; conn.line.width = Pt(w)
    no_shadow(conn)
    if label:
        mx, my = (x1+x2)/2 - 0.5, (y1+y2)/2 - 0.18
        lb = rect(s, mx, my, 1.0, 0.32, fill=ACCENT,
                  shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.5)
        boxtext(lb, [{'t': label, 's': 9.5, 'b': True, 'c': WHITE}])
line_between('unk', 'lead', ACCENT, 'follows', 2.5)
line_between('inf', 'unk', MGRAY)
line_between('col', 'lead', MGRAY)
for k, (x, y, t, col) in nodes.items():
    c = rect(s, x, y, 1.1, 0.7, fill=col, shape=MSO_SHAPE.OVAL)
    boxtext(c, [{'t': l, 's': 9.5, 'b': True, 'c': WHITE} for l in t.split('\n')])
# right explanation
txt(s, 6.5, 1.9, 6.2, 0.45, [{'t': '🔗  핵심 직관', 's': 14, 'b': True, 'c': NAVY}])
txt(s, 6.5, 2.35, 6.2, 1.6,
    [{'runs': [{'t': '“미상 표적”', 's': 12.5, 'b': True, 'c': GRAY},
               {'t': '은 위치만으론 선두 전차와 구별 불가.', 's': 12.5, 'c': DARK}], 'sa': 3},
     {'runs': [{'t': '하지만 ', 's': 12.5, 'c': DARK},
               {'t': '‘미군 전차를 follows(뒤따른다)’', 's': 12.5, 'b': True, 'c': ACCENT},
               {'t': '는 관계로 정체가 드러난다.', 's': 12.5, 'c': DARK}]}])
mech = [
    ('14종 관계 엣지', 'follows·firesAt·occupies·withdrawsFrom·reinforces …'),
    ('RGAT × 2층 (4헤드)', '중요한 이웃에 더 큰 가중치(attention)'),
    ('2-hop 전파', '이웃의 이웃까지 정보 도달'),
    ('잔차 연결 (h+Δ)', '자기 정체성 보존 + 맥락만 보강'),
]
y = 3.95
for t, d in mech:
    chip(s, 6.5, y, 3.1, t, BLUE, size=11, h=0.42)
    txt(s, 9.75, y, 3.0, 0.42, [{'t': d, 's': 10, 'c': DARK}], anchor=MSO_ANCHOR.MIDDLE)
    y += 0.55
res = rect(s, 0.7, 6.3, 11.9, 0.62, fill=RGBColor(0xEA,0xF3,0xEC),
           line=GREEN, line_w=1.3, shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.1)
boxtext(res, [{'runs': [
    {'t': '✔ 실험 효과:  ', 's': 12.5, 'b': True, 'c': GREEN},
    {'t': '한쪽이 UNKNOWN인 관측을 문맥으로 올바른 타입(CCF_Infantry)으로 식별해 병합 성공 (확신도 0.92)',
     's': 12.5, 'c': DARK}]}])

# ============================================================ SLIDE 10: CTA
s = new()
header(s, '3. 방법론  |  ⑤  CTA — 제약-인지 전이 어텐션', '모델의 두뇌: 6개 항을 “더해서” 설명가능하게 판정', 10)
# formula bar
fb = rect(s, 0.7, 1.3, 11.9, 0.7, fill=NAVY, shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.1)
boxtext(fb, [{'runs': [
    {'t': 'score = ', 's': 15, 'b': True, 'c': WHITE},
    {'t': 'w₀·의미유사도 ', 's': 14, 'b': True, 'c': RGBColor(0xF2,0xC2,0x8B)},
    {'t': '+ 시간 + 운동학 + 상태 + 관계 + 출처', 's': 14, 'c': WHITE},
    {'t': '     →  p = sigmoid(score)', 's': 13, 'b': True, 'c': RGBColor(0x9F,0xD8,0xB0)}]}])
terms = [
    ('의미 유사도', '두 벡터 코사인', '학습', BLUE),
    ('시간 (HARD)', '역행이면 −50 차단', '하드', RED),
    ('운동학', '필요속도>한계면 벌점', '학습', ACCENT),
    ('상태 호환', '교전→파괴 OK 등', '학습', STEEL),
    ('관계 겹침', '이웃 자카드', '학습', PURPLE),
    ('출처 편향', '출처쌍별 보정', '학습', GREEN),
]
x = 0.7
w6 = (11.9 - 5*0.16) / 6
for t, d, kind, col in terms:
    c = rect(s, x, 2.2, w6, 1.4, fill=WHITE, line=col, line_w=1.6,
             shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.06)
    txt(s, x+0.05, 2.3, w6-0.1, 1.25,
        [{'t': t, 's': 11.5, 'b': True, 'c': col, 'a': PP_ALIGN.CENTER, 'sa': 3},
         {'t': d, 's': 9.5, 'c': DARK, 'a': PP_ALIGN.CENTER, 'sa': 3},
         {'t': f'[{kind}]', 's': 9, 'b': True, 'c': (RED if kind=='하드' else GRAY),
          'a': PP_ALIGN.CENTER}])
    x += w6 + 0.16
# example: kinematics
ex = rect(s, 0.7, 3.85, 6.0, 2.5, fill=RGBColor(0xFB,0xEC,0xE6),
          shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.05)
txt(s, 0.9, 3.98, 5.6, 2.3,
    [{'t': '🏃  전장 예시 — 운동학 제약', 's': 13, 'b': True, 'c': ACCENT, 'sa': 5},
     {'runs': [{'t': '두 “보병” 관측이 의미상 비슷해도 ', 's': 12, 'c': DARK},
               {'t': '그 사이를 시속 344km(95 m/s)', 's': 12, 'b': True, 'c': RED},
               {'t': '로 달려야 한다면?', 's': 12, 'c': DARK}], 'sa': 4},
     {'t': '필요속도 = 거리 ÷ 시간', 's': 11.5, 'c': DARK, 'a': PP_ALIGN.CENTER, 'sa': 2},
     {'t': '가능속도 = v_max(타입,상태) + 위치오차', 's': 11.5, 'c': DARK, 'a': PP_ALIGN.CENTER, 'sa': 5},
     {'runs': [{'t': '⇒ ', 's': 12, 'b': True, 'c': ACCENT},
               {'t': '필요 > 가능 → 운동학 벌점 → 병합 차단', 's': 12, 'b': True, 'c': DARK}]}])
# why additive
ex2 = rect(s, 6.95, 3.85, 5.65, 2.5, fill=LGRAY,
           shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.05)
txt(s, 7.15, 3.98, 5.3, 2.3,
    [{'t': '💡  왜 “더하기” 구조인가 — 설명가능성', 's': 13, 'b': True, 'c': NAVY, 'sa': 5},
     {'t': '점수가 6개 항의 합이라, 어느 항 때문에 합쳐졌/막혔는지 분해해 볼 수 있다.',
      's': 12, 'c': DARK, 'sa': 5},
     {'runs': [{'t': '“시간”만 하드(−50): ', 's': 11.5, 'b': True, 'c': RED},
               {'t': '미래→과거 전이는 절대 불가.', 's': 11.5, 'c': DARK}], 'sa': 3},
     {'runs': [{'t': '나머지는 소프트 벌점: ', 's': 11.5, 'b': True, 'c': ACCENT},
               {'t': '센서 오차 여지를 남김.', 's': 11.5, 'c': DARK}], 'sa': 3},
     {'t': '→ 군사 의사결정의 “왜?”에 답할 수 있다.', 's': 11.5, 'it': True, 'c': GRAY}])

# ============================================================ SLIDE 11: heads + decoder
s = new()
header(s, '3. 방법론  |  ⑥  예측 헤드 · Sinkhorn 디코더', '쌍 판정 + 전체를 “부대 상자”로 군집화', 11)
# left: heads
txt(s, 0.7, 1.35, 6.0, 0.4, [{'t': '⑥a  예측 헤드 (쌍 단위)', 's': 14, 'b': True, 'c': GREEN}])
ph = rect(s, 0.7, 1.8, 5.6, 1.05, fill=WHITE, line=GREEN, line_w=1.5,
          shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.06)
txt(s, 0.85, 1.9, 5.3, 0.95,
    [{'t': 'PairHead — “이 둘은 같은가?”', 's': 12.5, 'b': True, 'c': DARK, 'sa': 3},
     {'t': '입력: 두 벡터 + |차이| + 곱 + CTA 6항 + dt/거리', 's': 11, 'c': GRAY}])
dh = rect(s, 0.7, 2.95, 5.6, 1.05, fill=WHITE, line=RED, line_w=1.5,
          shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.06)
txt(s, 0.85, 3.05, 5.3, 0.95,
    [{'t': 'DanglingHead — “매칭불가인가?”', 's': 12.5, 'b': True, 'c': DARK, 'sa': 3},
     {'t': '좋은 짝후보가 없으면 → 한쪽에만 존재할 가능성↑', 's': 11, 'c': GRAY}])
# left example
exL = rect(s, 0.7, 4.2, 5.6, 2.2, fill=RGBColor(0xEA,0xF3,0xEC),
           shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.05)
txt(s, 0.9, 4.33, 5.2, 2.0,
    [{'t': '⊘  전장 예시 — 매칭불가(dangling)', 's': 12.5, 'b': True, 'c': GREEN, 'sa': 5},
     {'t': 'KG-B만 잡은 “중공군 박격포대”는 KG-A에 대응 부대가 없다.', 's': 11.5, 'c': DARK, 'sa': 4},
     {'t': '→ 후보 점수가 모두 낮음 → ∅(null)로 분류 = 올바른 “매칭불가” 정탐.',
      's': 11.5, 'c': DARK}])
# right: decoder
txt(s, 6.7, 1.35, 6.0, 0.4, [{'t': '⑥b  Sinkhorn 슬롯 디코더 (전체 군집)', 's': 14, 'b': True, 'c': GREEN}])
# slots visual
slot_x = 6.9
for i, (lbl, col, obs) in enumerate([
        ('슬롯 1', BLUE, ['A:보병대대', 'B:미상보병']),
        ('슬롯 2', STEEL, ['A:전차', 'B:ROK tank']),
        ('∅ null', GRAY, ['B:박격포대'])]):
    box = rect(s, slot_x, 1.85 + i*0.95, 5.7, 0.82,
               fill=(RGBColor(0xF1,0xF4,0xF8)),
               line=col, line_w=1.5, shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.08)
    chip(s, slot_x+0.12, 1.97 + i*0.95, 1.1, lbl, col, size=10.5, h=0.32)
    ox = slot_x + 1.4
    for o in obs:
        oc = rect(s, ox, 1.97 + i*0.95, 1.55, 0.32, fill=WHITE, line=col, line_w=1.0,
                  shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.3)
        boxtext(oc, [{'t': o, 's': 9.5, 'b': True, 'c': col}])
        ox += 1.7
txt(s, 6.7, 4.75, 6.0, 1.7,
    [{'t': '🧩  핵심 아이디어', 's': 12.5, 'b': True, 'c': NAVY, 'sa': 4},
     {'runs': [{'t': '슬롯 = “빈 부대 상자” K개 + ∅ 1개. ', 's': 11.5, 'c': DARK},
               {'t': '같은 상자 = 같은 부대', 's': 11.5, 'b': True, 'c': GREEN},
               {'t': '.', 's': 11.5, 'c': DARK}], 'sa': 3},
     {'t': 'Sinkhorn: 한 상자에 다 몰리는 “붕괴”를 막는 균형 배정. ∅는 균형에서 빼 dangling을 자유 흡수.',
      's': 11.5, 'c': DARK, 'sa': 3},
     {'t': 'CTA 전이확률이 높은 관측끼리 같은 상자로 끌리게 주입 → ⑤와 ⑥b가 일관되게 정렬.',
      's': 11.5, 'c': DARK}])

# ============================================================ SLIDE 12: experiments dataset + quant
s = new()
header(s, '4. 실험 결과  |  EXPERIMENTS', '데이터셋 규모와 핵심 정량 지표', 12)
# dataset chips
dchips = [('1,005,101', '트리플'), ('75', '섹터 (50/12/13)'), ('51,310', '관측'),
          ('40', 'epoch'), ('128', 'd_model')]
x = 0.7
for v, l in dchips:
    c = rect(s, x, 1.3, 2.25, 0.85, fill=NAVY, shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.1)
    boxtext(c, [{'t': v, 's': 19, 'b': True, 'c': WHITE},
                {'t': l, 's': 10, 'c': RGBColor(0xCF,0xDD,0xEC)}])
    x += 2.4
# dataset distribution figure
if (FIG / 'dataset_distributions.png').exists():
    s.shapes.add_picture(str(FIG / 'dataset_distributions.png'),
                         Inches(0.7), Inches(2.35), width=Inches(7.0))
txt(s, 0.7, 5.65, 7.1, 1.5,
    [{'t': '데이터 분포 관찰', 's': 12.5, 'b': True, 'c': NAVY, 'sa': 3},
     {'runs': [{'t': '· 보병 편중: ', 's': 11.5, 'b': True, 'c': RED},
               {'t': 'CCF+ROK Infantry가 전체의 54% (클래스 불균형)', 's': 11.5, 'c': DARK}], 'sa': 2},
     {'t': '· 상태: 교전(Engaging)·사격(Firing) 위주 / 출처: KG-A·B 균형', 's': 11.5, 'c': DARK}])
# key metrics big
mt = rect(s, 8.0, 2.35, 4.6, 2.4, fill=LGRAY, shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.05)
txt(s, 8.2, 2.5, 4.3, 2.2,
    [{'t': '★ test 핵심 지표', 's': 13, 'b': True, 'c': NAVY, 'a': PP_ALIGN.CENTER, 'sa': 6},
     {'runs': [{'t': 'F1  ', 's': 14, 'b': True, 'c': DARK},
               {'t': '0.686', 's': 22, 'b': True, 'c': ACCENT}], 'a': PP_ALIGN.CENTER, 'sa': 4},
     {'runs': [{'t': 'Precision  ', 's': 13, 'c': DARK},
               {'t': '0.895', 's': 16, 'b': True, 'c': GREEN},
               {'t': '   Recall  ', 's': 13, 'c': DARK},
               {'t': '0.561', 's': 16, 'b': True, 'c': RED}], 'a': PP_ALIGN.CENTER, 'sa': 6}])
txt(s, 8.0, 4.95, 4.7, 2.0,
    [{'t': '해석', 's': 12.5, 'b': True, 'c': NAVY, 'sa': 3},
     {'runs': [{'t': '고정밀·저재현: ', 's': 11.5, 'b': True, 'c': ACCENT},
               {'t': '“확실하지 않으면 합치지 않는” 보수적 정책.', 's': 11.5, 'c': DARK}], 'sa': 3},
     {'t': '→ 오병합(다른 부대 오인)은 회피하나, 정합쌍을 절반 가까이 놓침(단편화 0.685).',
      's': 11.5, 'c': DARK}])

# ============================================================ SLIDE 13: experiments compare
s = new()
header(s, '4. 실험 결과  |  EXPERIMENTS', '기존 대비 향상 · 학습 수렴 · 섹터별 일관성', 13)
if (FIG / 'baseline_compare.png').exists():
    s.shapes.add_picture(str(FIG / 'baseline_compare.png'),
                         Inches(0.65), Inches(1.35), width=Inches(4.3))
if (FIG / 'training_curve.png').exists():
    s.shapes.add_picture(str(FIG / 'training_curve.png'),
                         Inches(5.1), Inches(1.35), width=Inches(4.2))
if (FIG / 'per_sector_metrics.png').exists():
    s.shapes.add_picture(str(FIG / 'per_sector_metrics.png'),
                         Inches(0.65), Inches(4.45), width=Inches(8.6))
# right notes
notes = rect(s, 9.5, 1.35, 3.25, 5.3, fill=LGRAY, shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.04)
txt(s, 9.68, 1.5, 2.95, 5.1,
    [{'t': '📈 읽는 법', 's': 13, 'b': True, 'c': NAVY, 'sa': 6},
     {'runs': [{'t': 'baseline 대비\n', 's': 11.5, 'b': True, 'c': GREEN},
               {'t': 'F1 0.52→0.69 (+32%)\nRecall 0.40→0.56 (+40%)', 's': 11, 'c': DARK}], 'sa': 8},
     {'runs': [{'t': '학습 수렴\n', 's': 11.5, 'b': True, 'c': BLUE},
               {'t': 'loss 17.6→1.63 안정 수렴.\ndev-F1는 30ep서 포화\n→ early-stop 여지', 's': 11, 'c': DARK}], 'sa': 8},
     {'runs': [{'t': '섹터별\n', 's': 11.5, 'b': True, 'c': ACCENT},
               {'t': 'Precision(빨강)은\n어디서나 높고 안정.\nRecall(초록)이 출렁\n→ 난이도=recall', 's': 11, 'c': DARK}], 'sa': 8},
     {'runs': [{'t': '일반화\n', 's': 11.5, 'b': True, 'c': PURPLE},
               {'t': 'test ≥ dev → 과적합\n징후 없음', 's': 11, 'c': DARK}]}])

# ============================================================ SLIDE 14: qualitative
s = new()
header(s, '4. 실험 결과  |  정성 분석', '실제 정합 사례 — 무엇을 잘하고 무엇을 틀리나', 14)
cases = [
    ('✔ 올바른 병합', GREEN,
     'A “중공군 침투조” ↔ B “미상 보병(UNKNOWN)”\n→ CCF_Infantry로 식별·병합 (확신도 0.92)',
     '이종 출처·UNKNOWN을 문맥으로 통합'),
    ('✘ 잘못된 병합', RED,
     'A “국군 보병 대대” ↔ B “국군 방어중대”\n(서로 다른 부대인데 병합)',
     '동일 타입 보병끼리 혼동 — 오병합의 71%'),
    ('✂ 단편화', ACCENT,
     'A “아군 보병” / B “아군 보병” (같은 실체)\n→ 확신 부족으로 미회수',
     'Recall 손실의 실체 = 보수적 abstain'),
    ('🛰 불가능 전이', PURPLE,
     'Approaching→Approaching, 필요속도 95.5 m/s\n(보병이 시속 344km)',
     '묶인 궤적의 24.6%가 물리적으로 불가능'),
]
positions = [(0.7, 1.4), (6.75, 1.4), (0.7, 4.05), (6.75, 4.05)]
for (t, col, body, note), (x, y) in zip(cases, positions):
    card = rect(s, x, y, 5.85, 2.4, fill=WHITE, line=col, line_w=2.0,
                shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.05)
    hd = rect(s, x, y, 5.85, 0.5, fill=col, shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.08)
    boxtext(hd, [{'t': t, 's': 13, 'b': True, 'c': WHITE}], align=PP_ALIGN.LEFT)
    hd.text_frame.margin_left = Pt(12)
    txt(s, x+0.2, y+0.62, 5.5, 1.7,
        [{'t': l, 's': 11.5, 'c': DARK, 'sa': 2} for l in body.split('\n')] +
        [{'runs': [{'t': '⇒ ', 's': 11, 'b': True, 'c': col},
                   {'t': note, 's': 11, 'b': True, 'c': GRAY}], 'sb': 4}])

# ============================================================ SLIDE 15: conclusion
s = new()
header(s, '5. 결론  |  CONCLUSION', '기여점 · 한계 · 향후 연구', 15)
# contributions
txt(s, 0.7, 1.3, 12, 0.4, [{'t': '✦  연구 기여점', 's': 14, 'b': True, 'c': GREEN}])
contribs = [
    ('설명가능 제약 융합', 'CTA 6항 가산식 — 의미+시간+운동학+상태+관계+출처를 분해가능하게 통합'),
    ('출처-인지 멀티모달', '출처 신뢰도 게이트로 센서별 강·약점 자동 보정'),
    ('정합불가 + 군집 동시', 'Dangling 헤드 + Sinkhorn 슬롯으로 매칭불가까지 처리'),
]
y = 1.75
for t, d in contribs:
    chip(s, 0.7, y, 3.0, t, GREEN, size=11, h=0.5)
    txt(s, 3.9, y, 8.7, 0.5, [{'t': d, 's': 11.5, 'c': DARK}], anchor=MSO_ANCHOR.MIDDLE)
    y += 0.6
# results recap chip row
rb = rect(s, 0.7, 3.6, 11.9, 0.6, fill=NAVY, shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.1)
boxtext(rb, [{'runs': [
    {'t': '실험 요약:  ', 's': 12.5, 'b': True, 'c': RGBColor(0xF2,0xC2,0x8B)},
    {'t': 'test F1 0.686 · Precision 0.895 · 기존 대비 F1 +32% / Recall +40% · 과적합 없음',
     's': 12.5, 'c': WHITE}]}])
# limitations & future side by side
lim = rect(s, 0.7, 4.45, 5.85, 2.25, fill=RGBColor(0xFB,0xEC,0xE6),
           shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.05)
txt(s, 0.9, 4.57, 5.5, 2.1,
    [{'t': '⚠  한계', 's': 13, 'b': True, 'c': RED, 'sa': 5},
     {'t': '· 저재현(Recall 0.56) — 단편화 0.685로 절반 가까이 놓침', 's': 11.5, 'c': DARK, 'sa': 3},
     {'t': '· 동일 타입 보병 혼동(오병합의 71%)', 's': 11.5, 'c': DARK, 'sa': 3},
     {'t': '· 불가능 전이 24.6% — CTA가 학습신호일 뿐 추론 시 미강제', 's': 11.5, 'c': DARK}])
fut = rect(s, 6.75, 4.45, 5.85, 2.25, fill=RGBColor(0xEA,0xF3,0xEC),
           shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.05)
txt(s, 6.95, 4.57, 5.5, 2.1,
    [{'t': '➜  향후 연구', 's': 13, 'b': True, 'c': GREEN, 'sa': 5},
     {'t': '· 임계값 재조정 — 저확신 오병합 제거로 P·R 동시 개선', 's': 11.5, 'c': DARK, 'sa': 3},
     {'t': '· CTA 제약을 추론 시 hard-gate화 → 불가능 전이 후처리 제거', 's': 11.5, 'c': DARK, 'sa': 3},
     {'t': '· 보병 hard-negative 대조학습 / UNKNOWN 전용 처리', 's': 11.5, 'c': DARK}])

# ============================================================ SLIDE 16: thanks
s = new()
rect(s, 0, 0, SW, SH, fill=NAVY)
rect(s, 0, 0, SW, 0.22, fill=ACCENT)
rect(s, 0, SH-0.22, SW, 0.22, fill=ACCENT)
txt(s, 0, 2.6, SW, 1.2,
    [{'t': '감사합니다', 's': 44, 'b': True, 'c': WHITE, 'a': PP_ALIGN.CENTER, 'sa': 6},
     {'t': 'Q & A', 's': 24, 'b': True, 'c': RGBColor(0xF2,0xC2,0x8B), 'a': PP_ALIGN.CENTER}])
txt(s, 0, 4.5, SW, 0.6,
    [{'t': 'CITA-Net · 백마고지(Hill 395) STKG 개체정합 실험', 's': 14,
      'c': RGBColor(0xCF,0xDD,0xEC), 'a': PP_ALIGN.CENTER}])

prs.save(str(OUT))
print('SAVED:', OUT)
print('slides:', len(prs.slides._sldIdLst))
